import streamlit as st
from docx import Document
import PyPDF2
import re
from pptx import Presentation
import json


ALLOWED_EXTENSIONS = ['txt', 'docx', 'pptx', 'pdf']

def clean_text(text):
    # Remove all escape characters
    cleaned_text = re.sub(r'[\n\r\t\f\v]', ' ', text)

    # Remove any other non-printable characters
    cleaned_text = re.sub(r'[^\x20-\x7E]', '', cleaned_text)

    # Replace multiple spaces with a single space
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text)

    # Strip leading and trailing whitespace
    cleaned_text = cleaned_text.strip()

    return cleaned_text

def extract_text_from_pdf(pdf_file_path):
    pdf_reader = PyPDF2.PdfReader(pdf_file_path)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return clean_text(text)

def extract_text_from_pptx(pptx_file_path):
    # Load the pptx file
    prs = Presentation(pptx_file_path)

    text = []

    # Loop through each slide
    for slide in prs.slides:
        # Loop through each shape in the slide 
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return clean_text("\n".join(text))

def extract_text_from_txt(txt_file_path):
    text = txt_file_path.read().decode('utf-8')
    return clean_text(text)

def extract_text_from_docx(docx_file_path):
    # Load the docx file
    doc = Document(docx_file_path)

    # Extract text from each paragraph in the document
    text = "\n".join([para.text for para in doc.paragraphs])
    return clean_text(text)


# Streamlit interface
st.title("Pre-Meeting Document Management")

with st.form(key = "form"):
    
    # Step 1: Add Discussion Points
    st.subheader("Add Discussion Points")
    discussion_points = st.text_area("Enter your discussion points (new point in new line)", height=200)
        
    # Step 2: Upload Meeting Documents
    st.subheader("Upload Relevant Documents")
    uploaded_files = st.file_uploader("Upload files (txt, docx, pptx, pdf)", type=ALLOWED_EXTENSIONS, accept_multiple_files=True)

    submitted = st.form_submit_button("Submit")
    if submitted:
        # store file text
        documents_content = ""
        if uploaded_files:
            for uploaded_file in uploaded_files:
                #uploaded_file = uploaded_files[i]
                filename = uploaded_file.name
                file_extension = filename.split('.')[-1].lower()

                # Extract text based on file type
                if file_extension == 'txt':
                    file_text = extract_text_from_txt(uploaded_file)
                elif file_extension == 'docx':
                    file_text = extract_text_from_docx(uploaded_file)
                elif file_extension == 'pptx':
                    file_text = extract_text_from_pptx(uploaded_file)
                elif file_extension == 'pdf':
                    file_text = extract_text_from_pdf(uploaded_file)
                else:
                    file_text = "Unsupported file format."

                # Store the file content
                documents_content = documents_content + file_text

        
        discussion_points = [x.strip() for x in discussion_points.split("\n")]
        
        
        with open("database/discussion_points.json", "w") as file:
            json.dump(discussion_points, file)

        with open("database/documents_content.txt", "w") as file:
            file.write(documents_content)



        



