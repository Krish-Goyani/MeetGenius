from docx import Document
import PyPDF2
import re
from pptx import Presentation
from  MeetGenius_logger import logger

"""
task1 : Pre-Meeting Document Management
- all docuemnts uploaded from users first of all are converted into text and than cleaned.
- alnog with that discussion points enterd from user stored in json formate for easy retrieval,
"""
def clean_text(text):
    # Remove all escape characters
    cleaned_text = re.sub(r'[\n\r\t\f\v]', ' ', text)

    # Remove any other non-printable characters
    cleaned_text = re.sub(r'[^\x20-\x7E]', '', cleaned_text)

    # Replace multiple spaces with a single space
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text)

    # Strip leading and trailing whitespace
    cleaned_text = cleaned_text.strip()

    return cleaned_text

def extract_text_from_pdf(pdf_file_path):
    #read pdf file and extract in to text

    pdf_reader = PyPDF2.PdfReader(pdf_file_path)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return clean_text(text)

def extract_text_from_pptx(pptx_file_path):
    #read pptx file and extract in to text
    # Load the pptx file
    prs = Presentation(pptx_file_path)

    text = []
    # Loop through each slide
    for slide in prs.slides:
        # Loop through each shape in the slide 
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return clean_text("\n".join(text))

def extract_text_from_txt(txt_file_path):
    #read txt file and extract in to text
    text = txt_file_path.read().decode('utf-8')
    return clean_text(text)

def extract_text_from_docx(docx_file_path):
    #read pdf docx and extract in to text

    # Load the docx file
    doc = Document(docx_file_path)

    # Extract text from each paragraph in the document
    text = "\n".join([para.text for para in doc.paragraphs])
    return clean_text(text)



def pre_process_documents(uploaded_files):
    """
    takes uploaded files than pass it in their extractor function according to their format
    and return accumulated content of all documents
    """

    logger.info("documents loading and text cleaning started")

    documents_content = ""

    for uploaded_file in uploaded_files:
        #uploaded_file = uploaded_files[i]
        filename = uploaded_file.name
        file_extension = filename.split('.')[-1].lower()

        # Extract text based on file type
        if file_extension == 'txt':
            file_text = extract_text_from_txt(uploaded_file)
        elif file_extension == 'docx':
            file_text = extract_text_from_docx(uploaded_file)
        elif file_extension == 'pptx':
            file_text = extract_text_from_pptx(uploaded_file)
        elif file_extension == 'pdf':
            file_text = extract_text_from_pdf(uploaded_file)
        else:
            file_text = "Unsupported file format."

        # Store the file content
        documents_content = documents_content + file_text

    logger.info("whole content accumulated and cleaned")

    return documents_content

